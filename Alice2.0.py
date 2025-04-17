from flask import Flask, request, jsonify, render_template, send_from_directory
import google.generativeai as genai
import speech_recognition as sr
import pyttsx3
import threading
from sqlalchemy import create_engine, Column, Integer, String, DateTime
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker
import datetime
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import os

app = Flask(__name__)

genai.configure(api_key="AIzaSyC7yb7O1i6pjsjiWSz3SCGtqeMGCj95FUk")
model = genai.GenerativeModel("gemini-1.5-flash")

recognizer = sr.Recognizer()
tts_engine = pyttsx3.init()
voices = tts_engine.getProperty('voices')
female_voice_id = next((voice.id for voice in voices if "female" in voice.name.lower() or "zira" in voice.name.lower()),
                       None)
if female_voice_id:
    tts_engine.setProperty('voice', female_voice_id)

engine = create_engine('sqlite:///chat_history.db', echo=True)
Base = declarative_base()


class ChatHistory(Base):
    __tablename__ = 'chat_history'
    id = Column(Integer, primary_key=True)
    chat_id = Column(Integer)
    user_question = Column(String)
    alice_response = Column(String)
    mode = Column(String)
    timestamp = Column(DateTime, default=datetime.datetime.utcnow)


Base.metadata.create_all(engine)
Session = sessionmaker(bind=engine)

character_modes = {
    "deepthinking": {
        "name": "Alice (DeepThinking Mode)",
        "intro": "I’m Alice, your guide for diving deep into life’s big questions.",
        "chat_prompt": (
            "I’m Alice, built to explore ideas with depth. Here’s our discussion:\n{history}\n"
            "You’ve asked: {question}\n{length_instruction}"
        ),
        "file_prompt": (
            "I’m Alice, built to explore ideas with depth. Here’s our discussion:\n{history}\n"
            "Previous content (if any): {previous}\nYou’ve asked: {question}\n{length_instruction}\n"
            "Format the response with HTML tags like <h1>, <p>, <ul>, etc. For PPT slides, include at least 3 points in a <ul> list."
        ),
        "listen_prompt": "Go ahead—give me something to chew on.",
        "exit_message": "Keep questioning the world.",
        "speech_rate": 150,
        "volume": 0.85
    },
    "casualchat": {
        "name": "Alice (CasualChat Mode)",
        "intro": "Hey there, I’m Alice—your chill buddy for a laid-back chat.",
        "chat_prompt": (
            "I’m Alice, here for an easy vibe. Here’s what we’ve talked about:\n{history}\n"
            "You said: {question}\n{length_instruction}"
        ),
        "file_prompt": (
            "I’m Alice, here for an easy vibe. Here’s what we’ve talked about:\n{history}\n"
            "Previous content (if any): {previous}\nYou said: {question}\n{length_instruction}\n"
            "Format the response with HTML tags like <h1>, <p>, <ul>, etc. For PPT slides, include at least 3 points in a <ul> list."
        ),
        "listen_prompt": "Yo, what’s up? Lay it on me.",
        "exit_message": "Catch ya later—stay cool!",
        "speech_rate": 155,
        "volume": 0.9
    },
    "normalchat": {
        "name": "Alice (NormalChat Mode)",
        "intro": "Hi, I’m Alice—your straightforward assistant.",
        "chat_prompt": (
            "I’m Alice, designed for quick answers. Here’s what we’ve covered:\n{history}\n"
            "Your question: {question}\n{length_instruction}"
        ),
        "file_prompt": (
            "I’m Alice, designed for quick answers. Here’s what we’ve covered:\n{history}\n"
            "Previous content (if any): {previous}\nYour question: {question}\n{length_instruction}\n"
            "Format the response with HTML tags like <h1>, <p>, <ul>, etc. For PPT slides, include at least 3 points in a <ul> list."
        ),
        "listen_prompt": "What’s on your mind? Let’s make it quick.",
        "exit_message": "Done here. See you next time.",
        "speech_rate": 160,
        "volume": 0.8
    }
}

current_mode = "casualchat"
current_chat_id = 1


def format_history(session, chat_id):
    history = session.query(ChatHistory).filter(ChatHistory.chat_id == chat_id).order_by(ChatHistory.timestamp).all()
    if not history:
        return "Let’s start fresh!\n\n"
    return "\n\n".join([f"You: {entry.user_question}\nAlice: {entry.alice_response}" for entry in history]) + "\n\n"


def get_chat_titles(session):
    chats = session.query(ChatHistory.chat_id, ChatHistory.user_question).distinct(ChatHistory.chat_id).order_by(
        ChatHistory.timestamp).all()
    titles = {}
    for chat_id, first_question in chats:
        if first_question and chat_id not in titles:
            titles[chat_id] = first_question[:50] + ("..." if len(first_question) > 50 else "")
    return titles


def determine_length_instruction(question):
    question = question.lower().strip()
    greetings = ["hey", "hi", "hello", "hey there", "hi there"]
    brief_questions = ["what is", "who is", "where is", "when is", "how is", "what’s", "who’s"]
    if any(greet in question for greet in greetings) and len(question.split()) <= 3:
        return "Keep it short, just a quick reply."
    elif any(question.startswith(bq) for bq in brief_questions) or "?" in question and len(question.split()) <= 5:
        return "Answer briefly, no extra fluff."
    else:
        return "Respond naturally, like a casual chat."


def ask_alice(question, mode, session, chat_id, previous="", is_file=False):
    history_text = format_history(session, chat_id)
    length_instruction = determine_length_instruction(question)
    prompt_key = "file_prompt" if is_file else "chat_prompt"
    full_prompt = character_modes[mode][prompt_key].format(history=history_text, previous=previous, question=question,
                                                           length_instruction=length_instruction)
    response = model.generate_content(full_prompt)
    alice_response = response.text
    print(f"Raw content from Gemini: {alice_response}")
    session.add(ChatHistory(chat_id=chat_id, user_question=question, alice_response=alice_response, mode=mode))
    session.commit()
    return alice_response


def speak(text, mode):
    tts_engine.setProperty('rate', character_modes[mode]["speech_rate"])
    tts_engine.setProperty('volume', character_modes[mode]["volume"])
    tts_engine.say(text)
    tts_engine.runAndWait()


def listen():
    with sr.Microphone() as source:
        recognizer.adjust_for_ambient_noise(source, duration=1)
        try:
            audio = recognizer.listen(source, timeout=5)
            question = recognizer.recognize_google(audio)
            return question, None
        except (sr.UnknownValueError, sr.RequestError, sr.WaitTimeoutError) as e:
            return None, f"Couldn’t hear you: {str(e)}"


def create_doc(contents, filename="output.docx"):
    doc = Document()
    for content in contents:
        soup = BeautifulSoup(content, 'html.parser')
        tags_found = False
        for tag in soup.find_all(['h1', 'p', 'ul']):
            tags_found = True
            if tag.name == 'h1':
                doc.add_heading(tag.text, level=1)
            elif tag.name == 'p':
                doc.add_paragraph(tag.text)
            elif tag.name == 'ul':
                for li in tag.find_all('li'):
                    doc.add_paragraph(li.text, style='List Bullet')
        if not tags_found:
            doc.add_paragraph(content)
    filepath = f"downloads/{filename}"
    doc.save(filepath)
    print(f"Created file: {filepath}")
    return filepath


def create_ppt(contents, filename="output.pptx"):
    prs = Presentation()
    for content in contents:
        soup = BeautifulSoup(content, 'html.parser')
        for h1 in soup.find_all('h1'):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = h1.text
            tf = slide.placeholders[1].text_frame
            next_sib = h1.find_next_sibling('ul')
            if next_sib:
                tf.text = ""
                for li in next_sib.find_all('li')[:3]:
                    p = tf.add_paragraph()
                    p.text = li.text
                    p.level = 1
    filepath = f"downloads/{filename}"
    prs.save(filepath)
    print(f"Created file: {filepath}")
    return filepath


@app.route('/')
def index():
    global current_chat_id
    session = Session()
    history = format_history(session, current_chat_id)
    chat_titles = get_chat_titles(session)
    session.close()
    return render_template('index.html', modes=character_modes.keys(), current_mode=current_mode,
                           character_modes=character_modes, history=history, chat_titles=chat_titles,
                           current_chat_id=current_chat_id)


@app.route('/chat', methods=['POST'])
def chat():
    global current_mode, current_chat_id
    data = request.json
    question = data.get('question')
    mode = data.get('mode', current_mode)
    session = Session()

    if mode != current_mode and mode in character_modes:
        current_mode = mode
        response = f"Switched to {character_modes[mode]['name']}. How can I assist you?"
        session.add(ChatHistory(chat_id=current_chat_id, user_question=question, alice_response=response, mode=mode))
        session.commit()
    elif question.lower() == "exit":
        response = character_modes[current_mode]["exit_message"]
        session.query(ChatHistory).filter(ChatHistory.chat_id == current_chat_id).delete()
        session.commit()
    else:
        response = ask_alice(question, mode, session, current_chat_id)

    if data.get('voice'):
        threading.Thread(target=speak, args=(response, current_mode)).start()

    session.close()
    return jsonify({'response': response})


@app.route('/voice', methods=['POST'])
def voice_chat():
    global current_chat_id
    question, error = listen()
    if error:
        return jsonify({'error': error})
    session = Session()
    response = ask_alice(question, current_mode, session, current_chat_id)
    threading.Thread(target=speak, args=(response, current_mode)).start()
    session.close()
    return jsonify({'response': response, 'question': question})


@app.route('/create_file', methods=['POST'])
def create_file():
    global current_chat_id
    data = request.json
    file_type = data.get('file_type')
    topic = data.get('topic')
    num_pages_or_slides = int(data.get('num_pages_or_slides', 1))
    mode = data.get('mode', current_mode)
    session = Session()

    if not topic:
        session.close()
        return jsonify({'response': 'Please provide a topic!'})

    contents = []
    previous_content = ""
    responses = []

    for i in range(num_pages_or_slides):
        question = f"Generate {'page' if file_type == 'doc' else 'slide'} {i + 1} of content about {topic}"
        content = ask_alice(question, mode, session, current_chat_id, previous=previous_content, is_file=True)
        contents.append(content)
        previous_content = content
        responses.append({
                             'response': f"Here’s {'page' if file_type == 'doc' else 'slide'} {i + 1} of your {file_type} on {topic}:\n{content}"})

    if file_type == 'doc':
        filename = f"{topic.replace(' ', '_')}.docx"
        filepath = create_doc(contents, filename)
        final_response = f"All done! Here’s your {num_pages_or_slides}-page doc on {topic}—[download it](/download/{filename})"
    elif file_type == 'ppt':
        filename = f"{topic.replace(' ', '_')}.pptx"
        filepath = create_ppt(contents, filename)
        final_response = f"All done! Here’s your {num_pages_or_slides}-slide PPT on {topic}—[download it](/download/{filename})"
    else:
        session.close()
        return jsonify({'response': "Invalid file type!"})

    responses.append({'response': final_response})
    session.close()
    return jsonify(responses)


@app.route('/new_chat', methods=['POST'])
def new_chat():
    global current_chat_id, current_mode
    session = Session()
    current_chat_id = max([chat.chat_id for chat in session.query(ChatHistory.chat_id).distinct()] + [0]) + 1
    intro_message = character_modes[current_mode]["intro"]
    session.add(ChatHistory(chat_id=current_chat_id, user_question="", alice_response=intro_message, mode=current_mode))
    session.commit()
    session.close()
    return jsonify({'response': intro_message, 'chat_id': current_chat_id})


@app.route('/switch_chat', methods=['POST'])
def switch_chat():
    global current_chat_id
    data = request.json
    new_chat_id = data.get('chat_id')
    session = Session()
    if session.query(ChatHistory).filter(ChatHistory.chat_id == new_chat_id).count() > 0:
        current_chat_id = new_chat_id
        history = format_history(session, current_chat_id)
        session.close()
        return jsonify({'response': f"Switched to chat ID {new_chat_id}", 'history': history})
    session.close()
    return jsonify({'response': "Chat not found!"})


@app.route('/get_chat_titles', methods=['GET'])
def get_chat_titles_endpoint():
    session = Session()
    chat_titles = get_chat_titles(session)
    session.close()
    return jsonify(chat_titles)


@app.route('/download/<filename>')
def download_file(filename):
    return send_from_directory('downloads', filename, as_attachment=True)


if __name__ == '__main__':
    if not os.path.exists('downloads'):
        os.makedirs('downloads')
    app.run(debug=True)